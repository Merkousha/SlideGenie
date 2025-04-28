import os
import subprocess
import tempfile
import shutil
import re
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

class MermaidGenerator:
    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
        # Check for mmdc in the system
        self.mmdc_path = self._find_mmdc_path()
        # Default size settings
        self.max_width = 12  # Maximum width (inches)
        self.max_height = 6  # Maximum height (inches)
        self.dpi = 96  # Default DPI for pixel to inch conversion

    def _find_mmdc_path(self):
        """
        Find mmdc path in the system
        """
        # Check for mmdc in PATH
        mmdc_path = shutil.which('mmdc')
        if mmdc_path:
            return mmdc_path
        
        # Check common npm installation paths
        common_paths = [
            os.path.expanduser('~/.npm-global/bin/mmdc'),
            os.path.expanduser('~/.nvm/versions/node/current/bin/mmdc'),
            os.path.expanduser('~/.nvm/versions/node/current/lib/node_modules/@mermaid-js/mermaid-cli/src/cli.js'),
            '/usr/local/bin/mmdc',
            '/usr/bin/mmdc',
            'C:\\Users\\%USERNAME%\\AppData\\Roaming\\npm\\mmdc.cmd',
            'C:\\Program Files\\nodejs\\mmdc.cmd'
        ]
        
        for path in common_paths:
            expanded_path = os.path.expandvars(path)
            if os.path.exists(expanded_path):
                return expanded_path
        
        return None

    def _calculate_dimensions_from_image(self, image_path):
        """
        Calculate appropriate dimensions for displaying the image in the slide based on actual image dimensions
        """
        try:
            with Image.open(image_path) as img:
                # Get image dimensions in pixels
                width_px, height_px = img.size
                
                # Convert pixels to inches
                width_inches = width_px / self.dpi
                height_inches = height_px / self.dpi
                
                # Calculate aspect ratio
                aspect_ratio = width_px / height_px
                
                # If width or height exceeds limits, scale the image
                if width_inches > self.max_width or height_inches > self.max_height:
                    if width_inches / self.max_width > height_inches / self.max_height:
                        # Width constraint
                        width_inches = self.max_width
                        height_inches = width_inches / aspect_ratio
                    else:
                        # Height constraint
                        height_inches = self.max_height
                        width_inches = height_inches * aspect_ratio
                
                return width_inches, height_inches
        except Exception as e:
            print(f"Error reading image dimensions: {str(e)}")
            # Return default dimensions in case of error
            return 8, 3

    def generate_image(self, mermaid_code, output_path):
        """
        Convert Mermaid code to image using mermaid-cli
        """
        if not self.mmdc_path:
            print("Error: mermaid-cli (mmdc) not found. Please install using 'npm install -g @mermaid-js/mermaid-cli'")
            return False
            
        # Create temporary file for Mermaid code
        mermaid_file = os.path.join(self.temp_dir, 'temp.mmd')
        with open(mermaid_file, 'w', encoding='utf-8') as f:
            f.write(mermaid_code)

        try:
            # Run mmdc command to convert to image
            subprocess.run([
                self.mmdc_path,
                '-i', mermaid_file,
                '-o', output_path,
                '-b', 'transparent'
            ], check=True)
            return True
        except subprocess.CalledProcessError as e:
            print(f"Error converting Mermaid to image: {str(e)}")
            return False
        except Exception as e:
            print(f"Unexpected error in Mermaid conversion: {str(e)}")
            return False
        finally:
            # Clean up temporary file
            if os.path.exists(mermaid_file):
                os.remove(mermaid_file)

    def add_mermaid_to_slide(self, slide, mermaid_code, width=None, height=None):
        """
        Add Mermaid diagram to the slide
        """
        # Create temporary path for storing the image
        temp_image = os.path.join(self.temp_dir, 'temp_mermaid.png')
        
        # Convert Mermaid to image
        if self.generate_image(mermaid_code, temp_image):
            # Calculate appropriate dimensions based on actual image dimensions
            if width is None or height is None:
                width_inches, height_inches = self._calculate_dimensions_from_image(temp_image)
                width = Inches(width_inches)
                height = Inches(height_inches)
            
            # Calculate center position
            slide_width = Inches(13.33)  # Standard PowerPoint slide width
            slide_height = Inches(7.5)   # Standard PowerPoint slide height
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            # Add image to slide
            slide.shapes.add_picture(temp_image, left, top, width, height)
            
            # Clean up temporary file
            if os.path.exists(temp_image):
                os.remove(temp_image)
            return True
        return False 