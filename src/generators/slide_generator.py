import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..config.settings import SLIDE_CONFIG, BASE_DIR
from .mermaid_generator import MermaidGenerator
from ..utils.logger import logger

class SlideGenerator:
    def __init__(self, output_path):
        self.output_path = output_path
        self.config = SLIDE_CONFIG
        self.prs = None
        self._initialize_presentation()
        self.mermaid_generator = MermaidGenerator()

    def _initialize_presentation(self):
        try:
            base_pptx_path = BASE_DIR / "Base.pptx"
            if not os.path.exists(base_pptx_path):
                raise FileNotFoundError(f"Base.pptx not found at {base_pptx_path}")
            os.system(f'copy "{base_pptx_path}" "{self.output_path}"')
            self.prs = Presentation(self.output_path)
            logger.info(f"Initialized presentation with base template at {base_pptx_path}")
        except Exception as e:
            logger.error(f"Error initializing presentation: {str(e)}")
            raise

    def add_title_slide(self, title, subtitle):
        try:
            title_slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = subtitle
            
            # Set title font size
            title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
            logger.info(f"Added title slide with title: {title} and subtitle: {subtitle}")
        except Exception as e:
            logger.error(f"Error adding title slide: {str(e)}")
            raise

    def add_content_slide(self, slide_data):
        try:
            content_slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(content_slide_layout)
            
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = slide_data["title"]
            content_shape.text = slide_data["content"]
            
            # Set font sizes
            title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(self.config["content_font_size"])
            logger.info(f"Added content slide with title: {slide_data['title']}")
        except Exception as e:
            logger.error(f"Error adding content slide: {str(e)}")
            raise

    def add_mermaid_diagram(self, slide, mermaid_code):
        try:
            result = self.mermaid_generator.add_mermaid_to_slide(slide, mermaid_code)
            if result:
                logger.info("Added Mermaid diagram to slide")
            else:
                logger.warning("Failed to add Mermaid diagram to slide")
            return result
        except Exception as e:
            logger.error(f"Error adding Mermaid diagram: {str(e)}")
            raise

    def add_code_examples(self, slide, code_examples):
        try:
            if not code_examples.strip():
                return
                
            code_text = f"""Code Examples:
            {code_examples}"""
            
            top = Inches(4.5)
            height = Inches(2)
            right = Inches(0.5)
            width = Inches(12)
            
            code_box = slide.shapes.add_textbox(right, top, width, height)
            code_frame = code_box.text_frame
            code_frame.text = code_text
            
            for paragraph in code_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.name = 'Consolas'
                paragraph.font.size = Pt(14)
            logger.info("Added code examples to slide")
        except Exception as e:
            logger.error(f"Error adding code examples: {str(e)}")
            raise

    def add_presenter_notes(self, slide, notes):
        try:
            if not notes:
                return
                
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes
            logger.info("Added presenter notes to slide")
        except Exception as e:
            logger.error(f"Error adding presenter notes: {str(e)}")
            raise

    def save(self):
        try:
            self.prs.save(self.output_path)
            logger.info(f"Presentation saved to {self.output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise
